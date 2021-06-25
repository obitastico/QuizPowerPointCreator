import json
import os
from os import PathLike
from os.path import abspath
from pathlib import Path
from random import shuffle
from time import sleep
from typing import Union, Optional

from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.shapes.autoshape import Shape
from psutil import process_iter
from win32com.client import Dispatch


def create_quiz_from_questions(questions):
    if "POWERPNT.EXE" in (p.name() for p in process_iter()):
        os.system("taskkill /F /IM POWERPNT.EXE")
        sleep(1)

    template_path = "./powerpoints/template.pptx"
    powerpoint_path = "./powerpoints/quiz.pptx"

    add_quiz_slides(template_path, powerpoint_path, 2, len(questions))

    prs = Presentation(powerpoint_path)

    for slide, question in zip(list(prs.slides)[1:], questions):
        slide.shapes.title.text = question["title"]
        shapes = list(filter(lambda s: "Abgerundetes Rechteck" in s.name and not s.text == "", slide.shapes))

        shuffle(shapes)

        right_shape = shapes[0]
        add_text(right_shape, question["richtig"])

        correct_shape = next(filter(lambda s: s.shape_id == 8, slide.shapes))
        add_text(correct_shape, right_shape.text)
        correct_shape.left = right_shape.left
        correct_shape.top = right_shape.top

        for i, shape in enumerate(shapes[1:]):
            add_text(shape, question["falsch"][i])

    prs.save(powerpoint_path)


def add_text(shape: Shape, text: str):
    shape.text += text
    shape.text_frame.fit_text(font_family="Century Gothic", max_size=40)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def load_questions(path: Union[str, Path, PathLike]):
    with open(path) as f:
        json_data = json.loads(f.read())
        try:
            return json_data["fragen"]
        except KeyError:
            raise Exception(f"Keine Fragen in {path} gefunden")


def add_quiz_slides(source_path: Union[str, Path, PathLike], destination_path: Union[str, Path, PathLike],
                    source_index: int, amount: Optional[int] = 1, *, read_only: Optional[bool] = True,
                    has_title: Optional[bool] = False, window: Optional[bool] = False):
    ppt_instance = Dispatch('PowerPoint.Application')

    prs = ppt_instance.Presentations.Open(abspath(source_path), read_only, has_title, window)

    prs.Slides(source_index).Copy()
    for i in range(amount-1):
        prs.Slides.Paste(Index=len(prs.Slides))

    prs.SaveAs(abspath(destination_path))
    prs.Close()

    ppt_instance.Quit()
    del ppt_instance
